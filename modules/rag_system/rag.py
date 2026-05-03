"""
RAG System Module (Retrieval Augmented Generation)

Upload PDF/PPTX -> Extract text -> Chunk -> Embed with google.genai (gemini-embedding-001)
-> Store numpy vectors -> Query: embed question -> cosine similarity search
-> top-K chunks -> Gemini LLM -> answer
"""

import os
import json
import numpy as np
from flask import request, jsonify
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from sklearn.metrics.pairwise import cosine_similarity
from google import genai
from langchain_google_genai import ChatGoogleGenerativeAI

load_dotenv()

ALLOWED_EXTENSIONS = {"pdf", "pptx", "ppt"}
CHUNK_SIZE      = 800
CHUNK_OVERLAP   = 100
TOP_K           = 4
PRIMARY_MODEL   = "gemini-2.5-flash-lite"
FALLBACK_MODEL  = "gemini-1.5-flash"
EMBEDDING_MODEL = "gemini-embedding-001"  # Uses google.genai client — tested & working


class RAGSystem:
    """
    Retrieval Augmented Generation for PDF/PPTX documents.

    Uses google.genai Client for embeddings (gemini-embedding-001)
    and cosine similarity (sklearn) for vector search.
    """

    def __init__(self, upload_folder: str, embeddings_folder: str, api_key: str):
        self.upload_folder     = upload_folder
        self.embeddings_folder = embeddings_folder
        self.api_key           = api_key
        os.makedirs(upload_folder,     exist_ok=True)
        os.makedirs(embeddings_folder, exist_ok=True)

        # Same API the user uses in their working embedding code
        self.genai_client = genai.Client(api_key=api_key)
        self.llm = ChatGoogleGenerativeAI(
            model=PRIMARY_MODEL, google_api_key=api_key, temperature=0)

    def allowed_file(self, filename: str) -> bool:
        """Returns True if file is PDF, PPTX, or PPT."""
        return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

    def embed_texts(self, texts: list) -> np.ndarray:
        """Embed a batch of text strings using google.genai client."""
        result = self.genai_client.models.embed_content(
            model=EMBEDDING_MODEL, contents=texts)
        return np.array([e.values for e in result.embeddings], dtype="float32")

    def embed_query(self, query: str) -> np.ndarray:
        """Embed a single query string."""
        result = self.genai_client.models.embed_content(
            model=EMBEDDING_MODEL, contents=query)
        return np.array([result.embeddings[0].values], dtype="float32")

    def extract_text(self, filepath: str, extension: str) -> str:
        """Extracts all readable text from a PDF or PPTX file."""
        text = ""
        if extension == "pdf":
            from pypdf import PdfReader
            for page in PdfReader(filepath).pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
        elif extension in ("pptx", "ppt"):
            from pptx import Presentation
            for slide in Presentation(filepath).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        return text

    def chunk_text(self, text: str) -> list:
        """Splits text into overlapping chunks."""
        chunks, start = [], 0
        while start < len(text):
            chunk = text[start:start + CHUNK_SIZE]
            if chunk.strip():
                chunks.append(chunk)
            start += CHUNK_SIZE - CHUNK_OVERLAP
        return chunks

    def save_index(self, doc_name: str, embeddings: np.ndarray, chunks: list):
        """Saves numpy embeddings + chunks JSON to disk."""
        np.save(os.path.join(self.embeddings_folder, f"{doc_name}.npy"), embeddings)
        with open(os.path.join(self.embeddings_folder, f"{doc_name}.json"), "w") as f:
            json.dump(chunks, f)

    def load_index(self, doc_name: str) -> tuple:
        """Loads saved embeddings and chunks from disk."""
        ep = os.path.join(self.embeddings_folder, f"{doc_name}.npy")
        cp = os.path.join(self.embeddings_folder, f"{doc_name}.json")
        if not os.path.exists(ep) or not os.path.exists(cp):
            raise FileNotFoundError(f"Document '{doc_name}' not indexed. Upload it first.")
        with open(cp) as f:
            chunks = json.load(f)
        return np.load(ep), chunks

    def similarity_search(self, question: str, doc_name: str) -> list:
        """Finds the top-K most relevant chunks using cosine similarity."""
        doc_embeddings, chunks = self.load_index(doc_name)
        q_emb = self.embed_query(question)
        scores = cosine_similarity(q_emb, doc_embeddings)[0]
        top_indices = np.argsort(scores)[::-1][:TOP_K]
        return [chunks[i] for i in top_indices if i < len(chunks)]

    def answer_question(self, question: str, doc_name: str) -> dict:
        """Full RAG pipeline: search -> answer with Gemini LLM."""
        chunks  = self.similarity_search(question, doc_name)
        context = "\n\n---\n\n".join(chunks)
        prompt  = f"""You are a helpful assistant. Answer ONLY using the context below.
If the answer is not in the context, say "I couldn't find that information in the document."

CONTEXT:
{context}

QUESTION: {question}

ANSWER:"""
        try:
            answer = self.llm.invoke(prompt).content
        except Exception:
            self.llm = ChatGoogleGenerativeAI(
                model=FALLBACK_MODEL, google_api_key=self.api_key, temperature=0)
            answer = self.llm.invoke(prompt).content
        return {"answer": answer, "source_chunks": chunks}

    def list_documents(self) -> list:
        """Returns names of all indexed documents (.npy files)."""
        return [f.replace(".npy", "")
                for f in os.listdir(self.embeddings_folder) if f.endswith(".npy")]

    def process_upload(self, filepath: str, filename: str) -> dict:
        """Full upload pipeline: extract -> chunk -> embed -> save."""
        ext      = filename.rsplit(".", 1)[1].lower()
        doc_name = filename.rsplit(".", 1)[0]
        try:
            text = self.extract_text(filepath, ext)
            if not text.strip():
                return {"success": False, "error": "No text could be extracted from the file."}
            chunks     = self.chunk_text(text)
            embeddings = self.embed_texts(chunks)
            self.save_index(doc_name, embeddings, chunks)
            return {"success": True, "doc_name": doc_name, "chunks_count": len(chunks)}
        except Exception as e:
            from modules.system_logger import SystemLogger
            SystemLogger.log("ERROR", "RAGSystem", f"Upload failed for '{filename}': {e}")
            return {"success": False, "error": str(e)}


class RAGRoutes:
    """Registers /rag/* Flask routes."""

    def __init__(self, app, rag: RAGSystem):
        self.app = app
        self.rag = rag

    def register(self):
        self.app.add_url_rule("/rag/upload",    "rag_upload",    self.upload,    methods=["POST"])
        self.app.add_url_rule("/rag/query",     "rag_query",     self.query,     methods=["POST"])
        self.app.add_url_rule("/rag/documents", "rag_documents", self.documents, methods=["GET"])

    def upload(self):
        """POST /rag/upload — Accepts PDF/PPTX, builds embedding index."""
        if "file" not in request.files:
            return jsonify({"error": "No file. Use key 'file'."}), 400
        file = request.files["file"]
        if not self.rag.allowed_file(file.filename):
            return jsonify({"error": "Only PDF and PPTX files are supported for RAG."}), 400
        filename = secure_filename(file.filename)
        filepath = os.path.join(self.rag.upload_folder, filename)
        file.save(filepath)
        result = self.rag.process_upload(filepath, filename)
        if not result["success"]:
            return jsonify(result), 500
        return jsonify({
            "success":      True,
            "doc_name":     result["doc_name"],
            "chunks_count": result["chunks_count"],
            "message":      f"Indexed {result['chunks_count']} chunks from '{filename}'."
        })

    def query(self):
        """POST /rag/query — Ask question from indexed document."""
        data = request.get_json()
        if not data or "question" not in data or "document" not in data:
            return jsonify({"error": "Provide 'question' and 'document'."}), 400
        try:
            result = self.rag.answer_question(data["question"], data["document"])
            return jsonify({
                "question":     data["question"],
                "document":     data["document"],
                "answer":       result["answer"],
                "source_chunks": result["source_chunks"]
            })
        except FileNotFoundError as e:
            return jsonify({"error": str(e)}), 404
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    def documents(self):
        """GET /rag/documents — list all indexed document names."""
        docs = self.rag.list_documents()
        return jsonify({"documents": docs, "count": len(docs)})
